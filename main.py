import logging

from pptx import Presentation
from pptx.shapes.placeholder import SlidePlaceholder
from pptx.util import Inches

# Generate Loggers
time_format = "%Y-%m-%d %H:%M:%S"
formatter = logging.Formatter(fmt='%(asctime)s - %(levelname)s - %(message)s', datefmt=time_format)

# Create a logger and set the custom formatter
logger = logging.getLogger("logger")
logger.setLevel(logging.INFO)
handler = logging.StreamHandler()

handler.setFormatter(formatter)
logger.addHandler(handler)

def WriteText(texto:str) -> bool:
    """
    Funcion para Escribir en Texto con diapositiva ya creada
    - Freddy Nava 13/12/2024
    :param texto:
    :return: Bool
    """
    try:
        nombre_file = "trial.pptx" # Nombre del Archivo
        logger.info(".......................................................................")
        logger.info("# Inicio del Script #")
        logger.info(".......................................................................")

        # Abrir presentancion Existente
        prs = Presentation(nombre_file)
        slide = prs.slides[0]

        # Recorro todas las figuras en las diapositivas
        # Si es tipo Place Holder muestro el texto de resto lo omito
        logger.info("Recorrer Elementos de la Diapositiva: ")
        for index,shape in enumerate(slide.shapes):
            logger.info(f"ID_ARRAY = {index} Forma ID: {shape.shape_id}, Nombre: {shape.name}, "
                        f"Value: {shape.text if isinstance(shape, SlidePlaceholder) else "Not Value Text"}")
        #
        logger.info(".......................................................................")
        title = slide.shapes[0] # Id 0 del Titulo en Especifico
        logger.info("Texto asignado a la Diapositiva...")
        logger.info(f"Texto en la diapositiva: {title.text}")
        # --------------------------------------------------------------------
        logger.info(".......................................................................")
        logger.info("Asignando Nuevo Texto...........")
        title.text = texto.upper()
        logger.info(f"Texto nuevo asignado: {title.text}")
        prs.save(nombre_file)
        return True
    except Exception as e:
        logger.error(e.args[0])
        return False

# Inicio del Programa
if __name__ == "__main__":
    if WriteText("Gracias!"):
        logger.info("Se Escribio Correctamente en el Texto")
    else:
        logger.error("Error al Escribir en el Texto")
